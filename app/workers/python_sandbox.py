#!/usr/bin/env python3
"""Restricted Python data worker for WES.

It intentionally exposes a small set of deterministic operations instead of
executing model-generated Python. This keeps the production worker useful for
analysis while preventing shell, network and arbitrary filesystem access.
"""

from __future__ import annotations

import ast
import csv
import hashlib
import html
import json
import math
import os
import statistics
import sys
from datetime import date, datetime
from pathlib import Path

MAX_ROWS = 50_000
MAX_TEXT = 500_000
MAX_COLUMNS = 100
MAX_SHEETS = 20
MAX_PAGES = 200
MAX_UNIQUE = 5_000
MAX_NUMERIC_SAMPLE = 10_000


def apply_limits() -> None:
    try:
        import resource
        resource.setrlimit(resource.RLIMIT_CPU, (8, 8))
        resource.setrlimit(resource.RLIMIT_AS, (256 * 1024 * 1024, 256 * 1024 * 1024))
        resource.setrlimit(resource.RLIMIT_FSIZE, (20 * 1024 * 1024, 20 * 1024 * 1024))
        resource.setrlimit(resource.RLIMIT_NOFILE, (32, 32))
    except (ImportError, ValueError, OSError):
        pass


def disable_network() -> None:
    import socket

    def blocked(*_args, **_kwargs):
        raise PermissionError("Accesso di rete non consentito nel worker Python")

    original_socket = socket.socket

    class RestrictedSocket(original_socket):
        """A real socket type so libraries can subclass it, without network I/O."""

        def connect(self, *_args, **_kwargs):
            return blocked()

        def connect_ex(self, *_args, **_kwargs):
            return blocked()

        def sendto(self, *_args, **_kwargs):
            return blocked()

    socket.socket = RestrictedSocket
    socket.create_connection = blocked


ROOT = Path(os.environ.get("AGENT_WORK_ROOT", "/tmp/wes-agent-workspaces")).resolve()


def safe_path(workspace: Path, relative: str) -> Path:
    candidate = (workspace / relative).resolve()
    if candidate != workspace and workspace not in candidate.parents:
        raise ValueError("Percorso fuori dal workspace")
    if candidate != ROOT and ROOT not in candidate.parents:
        raise ValueError("Workspace non autorizzato")
    return candidate


ALLOWED_BINARY_OPERATORS = {
    ast.Add: lambda a, b: a + b,
    ast.Sub: lambda a, b: a - b,
    ast.Mult: lambda a, b: a * b,
    ast.Div: lambda a, b: a / b,
    ast.FloorDiv: lambda a, b: a // b,
    ast.Mod: lambda a, b: a % b,
    ast.Pow: lambda a, b: a**b,
}
ALLOWED_UNARY_OPERATORS = {ast.UAdd: lambda a: a, ast.USub: lambda a: -a}
ALLOWED_FUNCTIONS = {
    "abs": abs, "round": round, "min": min, "max": max,
    "sqrt": math.sqrt, "log": math.log, "ceil": math.ceil, "floor": math.floor,
}


def evaluate_expression(expression: str) -> float:
    if len(expression) > 500:
        raise ValueError("Espressione troppo lunga")
    tree = ast.parse(expression, mode="eval")

    def evaluate(node):
        if isinstance(node, ast.Expression):
            return evaluate(node.body)
        if isinstance(node, ast.Constant) and isinstance(node.value, (int, float)):
            return node.value
        if isinstance(node, ast.BinOp) and type(node.op) in ALLOWED_BINARY_OPERATORS:
            left, right = evaluate(node.left), evaluate(node.right)
            if isinstance(node.op, ast.Pow) and abs(right) > 20:
                raise ValueError("Esponente non consentito")
            return ALLOWED_BINARY_OPERATORS[type(node.op)](left, right)
        if isinstance(node, ast.UnaryOp) and type(node.op) in ALLOWED_UNARY_OPERATORS:
            return ALLOWED_UNARY_OPERATORS[type(node.op)](evaluate(node.operand))
        if isinstance(node, ast.Call) and isinstance(node.func, ast.Name) and node.func.id in ALLOWED_FUNCTIONS:
            return ALLOWED_FUNCTIONS[node.func.id](*(evaluate(arg) for arg in node.args))
        raise ValueError("Operazione non consentita")

    result = evaluate(tree)
    if not isinstance(result, (int, float)) or not math.isfinite(result):
        raise ValueError("Risultato non valido")
    return result


def numeric(value):
    if isinstance(value, bool):
        return None
    if isinstance(value, (int, float)):
        result = float(value)
        return result if math.isfinite(result) and abs(result) <= 1e100 else None
    normalized = str(value).strip().replace(" ", "").replace(",", ".")
    try:
        result = float(normalized)
        return result if math.isfinite(result) and abs(result) <= 1e100 else None
    except ValueError:
        return None


def normalized_headers(values) -> list[str]:
    if len(values) > MAX_COLUMNS:
        raise ValueError(f"Il file supera il limite di {MAX_COLUMNS} colonne")
    headers, used = [], set()
    for index, value in enumerate(values):
        base = str(value or f"colonna_{index + 1}").strip()[:120] or f"colonna_{index + 1}"
        candidate, suffix = base, 2
        while candidate in used:
            candidate = f"{base}_{suffix}"
            suffix += 1
        used.add(candidate)
        headers.append(candidate)
    return headers


def new_column_state() -> dict:
    return {
        "missing": 0, "unique": set(), "unique_truncated": False,
        "numeric_count": 0, "numeric_sum": 0.0, "numeric_min": None,
        "numeric_max": None, "numeric_sample": [],
    }


def track_value(state: dict, value) -> None:
    if value is None or str(value).strip() == "":
        state["missing"] += 1
        return
    display = value.isoformat() if isinstance(value, (date, datetime)) else str(value).strip()
    if len(state["unique"]) < MAX_UNIQUE:
        state["unique"].add(display[:200])
    elif display[:200] not in state["unique"]:
        state["unique_truncated"] = True
    number = numeric(value)
    if number is None:
        return
    state["numeric_count"] += 1
    state["numeric_sum"] += number
    state["numeric_min"] = number if state["numeric_min"] is None else min(state["numeric_min"], number)
    state["numeric_max"] = number if state["numeric_max"] is None else max(state["numeric_max"], number)
    if len(state["numeric_sample"]) < MAX_NUMERIC_SAMPLE:
        state["numeric_sample"].append(number)


def summarize_columns(columns: dict, row_count: int) -> dict:
    fields = {}
    for header, data in columns.items():
        field = {
            "missing": data["missing"], "unique": len(data["unique"]),
            "unique_truncated": data["unique_truncated"],
        }
        populated = max(0, row_count - data["missing"])
        if data["numeric_count"] and data["numeric_count"] >= max(1, populated) * 0.8:
            sample = data["numeric_sample"]
            field.update({
                "type": "number", "min": data["numeric_min"], "max": data["numeric_max"],
                "mean": data["numeric_sum"] / data["numeric_count"],
                "median": statistics.median(sample),
                "median_is_sampled": len(sample) < data["numeric_count"],
            })
        else:
            field["type"] = "text"
            field["examples"] = sorted(data["unique"])[:5]
        fields[header] = field
    return fields


def analyze_csv(workspace: Path, payload: dict) -> dict:
    source = safe_path(workspace, str(payload.get("file", "")))
    if not source.is_file() or source.suffix.lower() not in {".csv", ".tsv"}:
        raise ValueError("Serve un file CSV o TSV valido")
    delimiter = "\t" if source.suffix.lower() == ".tsv" else ","
    with source.open("r", encoding="utf-8-sig", errors="replace", newline="") as handle:
        sample = handle.read(8192)
        handle.seek(0)
        try:
            delimiter = csv.Sniffer().sniff(sample, delimiters=",;\t|").delimiter
        except csv.Error:
            pass
        reader = csv.reader(handle, delimiter=delimiter)
        try:
            headers = normalized_headers(next(reader))
        except StopIteration:
            raise ValueError("Il file CSV è vuoto")
        columns = {header: new_column_state() for header in headers}
        row_count = 0
        for row in reader:
            row_count += 1
            if row_count > MAX_ROWS:
                raise ValueError(f"Il file supera il limite di {MAX_ROWS} righe")
            for index, header in enumerate(headers):
                track_value(columns[header], row[index] if index < len(row) else None)

    summary = {"rows": row_count, "columns": len(headers), "fields": summarize_columns(columns, row_count)}

    output_name = "analisi-dati.json"
    output = safe_path(workspace, output_name)
    output.write_text(json.dumps(summary, ensure_ascii=False, indent=2), encoding="utf-8")
    return {"summary": summary, "artifacts": [{"name": output_name, "path": str(output), "type": "application/json"}]}


def analyze_spreadsheet(workspace: Path, payload: dict) -> dict:
    from openpyxl import load_workbook

    source = safe_path(workspace, str(payload.get("file", "")))
    if not source.is_file() or source.suffix.lower() != ".xlsx":
        raise ValueError("Serve un file XLSX valido")
    workbook = load_workbook(source, read_only=True, data_only=True, keep_links=False)
    if len(workbook.sheetnames) > MAX_SHEETS:
        raise ValueError(f"Il file supera il limite di {MAX_SHEETS} fogli")
    result = {"sheets": {}}
    total_rows = 0
    for sheet_name in workbook.sheetnames:
        sheet = workbook[sheet_name]
        rows = sheet.iter_rows(values_only=True)
        try:
            headers = normalized_headers(list(next(rows)))
        except StopIteration:
            result["sheets"][sheet_name] = {"rows": 0, "columns": 0, "fields": {}}
            continue
        columns = {header: new_column_state() for header in headers}
        row_count = 0
        for row in rows:
            row_count += 1
            total_rows += 1
            if total_rows > MAX_ROWS:
                raise ValueError(f"Il file supera il limite complessivo di {MAX_ROWS} righe")
            for index, header in enumerate(headers):
                track_value(columns[header], row[index] if index < len(row) else None)
        result["sheets"][sheet_name] = {
            "rows": row_count, "columns": len(headers), "fields": summarize_columns(columns, row_count)
        }
    workbook.close()
    output_name = "analisi-excel.json"
    output = safe_path(workspace, output_name)
    output.write_text(json.dumps(result, ensure_ascii=False, indent=2), encoding="utf-8")
    return {"summary": result, "artifacts": [{"name": output_name, "path": str(output), "type": "application/json"}]}


def inspect_text(workspace: Path, payload: dict) -> dict:
    source = safe_path(workspace, str(payload.get("file", "")))
    if not source.is_file() or source.suffix.lower() not in {".txt", ".md", ".json"}:
        raise ValueError("Formato testuale non supportato")
    content = source.read_text(encoding="utf-8", errors="replace")[:MAX_TEXT]
    words = content.split()
    return {
        "characters": len(content), "words": len(words),
        "lines": content.count("\n") + 1,
        "sha256": hashlib.sha256(content.encode()).hexdigest(),
        "preview": content[:2000],
        "artifacts": [],
    }


def analyze_document(workspace: Path, payload: dict) -> dict:
    source = safe_path(workspace, str(payload.get("file", "")))
    suffix = source.suffix.lower()
    if not source.is_file() or suffix not in {".pdf", ".docx", ".pptx"}:
        raise ValueError("Formato documento non supportato")
    chunks, details, character_count = [], {"format": suffix.lstrip(".")}, 0

    def append_text(value) -> bool:
        nonlocal character_count
        remaining = MAX_TEXT - character_count
        if remaining <= 0:
            return False
        piece = str(value or "")[:remaining]
        chunks.append(piece)
        character_count += len(piece) + 1
        return character_count < MAX_TEXT

    if suffix == ".pdf":
        from pypdf import PdfReader
        reader = PdfReader(str(source), strict=False)
        if reader.is_encrypted and not reader.decrypt(""):
            raise ValueError("Il PDF è protetto da password")
        if len(reader.pages) > MAX_PAGES:
            raise ValueError(f"Il PDF supera il limite di {MAX_PAGES} pagine")
        details["pages"] = len(reader.pages)
        for page in reader.pages:
            if not append_text(page.extract_text() or ""):
                break
    elif suffix == ".docx":
        from docx import Document
        document = Document(str(source))
        for paragraph in document.paragraphs:
            if not append_text(paragraph.text):
                break
        for table in document.tables:
            for row in table.rows:
                if not append_text(" | ".join(cell.text for cell in row.cells)):
                    break
            if character_count >= MAX_TEXT:
                break
        details["paragraphs"] = len(document.paragraphs)
        details["tables"] = len(document.tables)
    else:
        from pptx import Presentation
        presentation = Presentation(str(source))
        if len(presentation.slides) > MAX_PAGES:
            raise ValueError(f"La presentazione supera il limite di {MAX_PAGES} slide")
        details["slides"] = len(presentation.slides)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False) and not append_text(shape.text):
                    break
            if character_count >= MAX_TEXT:
                break
    content = "\n".join(chunks)[:MAX_TEXT]
    output_name = "testo-estratto.txt"
    output = safe_path(workspace, output_name)
    output.write_text(content, encoding="utf-8")
    summary = {
        **details, "characters": len(content), "words": len(content.split()),
        "truncated": character_count >= MAX_TEXT, "preview": content[:4_000],
    }
    return {"summary": summary, "artifacts": [{"name": output_name, "path": str(output), "type": "text/plain"}]}


def inspect_image(workspace: Path, payload: dict) -> dict:
    from PIL import Image
    Image.MAX_IMAGE_PIXELS = 40_000_000
    source = safe_path(workspace, str(payload.get("file", "")))
    if not source.is_file() or source.suffix.lower() not in {".png", ".jpg", ".jpeg", ".webp"}:
        raise ValueError("Formato immagine non supportato")
    with Image.open(source) as image:
        if image.width * image.height > Image.MAX_IMAGE_PIXELS:
            raise ValueError("L’immagine supera il limite di 40 megapixel")
        image.verify()
    with Image.open(source) as image:
        return {
            "summary": {
                "format": image.format, "width": image.width, "height": image.height,
                "mode": image.mode, "frames": int(getattr(image, "n_frames", 1)),
                "note": "Sono stati verificati file e metadati; il riconoscimento visivo non è attivo in questo worker.",
            },
            "artifacts": [],
        }


def add_pdf_text(story, styles, value: str, style_name: str = "BodyText") -> None:
    from reportlab.platypus import Paragraph, Spacer
    clean = str(value).replace("\x00", "").strip()
    if not clean:
        return
    for start in range(0, len(clean), 2_000):
        story.append(Paragraph(html.escape(clean[start:start + 2_000]).replace("\n", "<br/>"), styles[style_name]))
        story.append(Spacer(1, 7))


def build_pdf(output: Path, title: str, sections: list) -> None:
    from reportlab.lib.enums import TA_LEFT
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import mm
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    from reportlab.platypus import SimpleDocTemplate

    font_path = Path("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf")
    font_name = "Helvetica"
    if font_path.is_file():
        pdfmetrics.registerFont(TTFont("WESDejaVu", str(font_path)))
        font_name = "WESDejaVu"
    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle(name="WESBody", parent=styles["BodyText"], fontName=font_name, fontSize=9.5, leading=14, alignment=TA_LEFT, textColor="#273449"))
    styles.add(ParagraphStyle(name="WESTitle", parent=styles["Title"], fontName=font_name, fontSize=22, leading=28, textColor="#171329", spaceAfter=14))
    styles.add(ParagraphStyle(name="WESHeading", parent=styles["Heading2"], fontName=font_name, fontSize=14, leading=18, textColor="#5B4CF0", spaceBefore=10, spaceAfter=6))
    document = SimpleDocTemplate(str(output), pagesize=A4, rightMargin=18 * mm, leftMargin=18 * mm, topMargin=18 * mm, bottomMargin=18 * mm, title=title, author="WES Autonomous Intelligence")
    story = []
    add_pdf_text(story, styles, title, "WESTitle")
    for section in sections:
        if not isinstance(section, dict):
            continue
        add_pdf_text(story, styles, str(section.get("title") or "Sezione")[:160], "WESHeading")
        add_pdf_text(story, styles, str(section.get("content") or "")[:50_000], "WESBody")
    document.build(story)


def create_report(workspace: Path, payload: dict) -> dict:
    title = str(payload.get("title") or "Report WES").strip()[:160]
    sections = payload.get("sections") or []
    if not isinstance(sections, list) or len(sections) > 40:
        raise ValueError("Sezioni report non valide")
    lines = [f"# {title}", "", "Generato nel workspace protetto WES.", ""]
    for section in sections:
        if not isinstance(section, dict):
            continue
        heading = str(section.get("title") or "Sezione").strip()[:160]
        body = str(section.get("content") or "").strip()[:200_000]
        lines.extend([f"## {heading}", "", body, ""])
    output_name = "report-wes.md"
    output = safe_path(workspace, output_name)
    output.write_text("\n".join(lines), encoding="utf-8")
    pdf_name = "report-wes.pdf"
    pdf_output = safe_path(workspace, pdf_name)
    build_pdf(pdf_output, title, sections)
    return {"artifacts": [
        {"name": output_name, "path": str(output), "type": "text/markdown"},
        {"name": pdf_name, "path": str(pdf_output), "type": "application/pdf"},
    ]}


OPERATIONS = {
    "calculate": lambda workspace, payload: {"result": evaluate_expression(str(payload.get("expression", ""))), "artifacts": []},
    "analyze_csv": analyze_csv,
    "analyze_spreadsheet": analyze_spreadsheet,
    "inspect_text": inspect_text,
    "analyze_document": analyze_document,
    "inspect_image": inspect_image,
    "create_report": create_report,
}


def main() -> None:
    apply_limits()
    disable_network()
    request = json.loads(sys.stdin.read(MAX_TEXT))
    workspace = Path(str(request.get("workspace", ""))).resolve()
    if workspace != ROOT and ROOT not in workspace.parents:
        raise ValueError("Workspace non autorizzato")
    operation = str(request.get("operation", ""))
    if operation not in OPERATIONS:
        raise ValueError("Operazione Python non consentita")
    result = OPERATIONS[operation](workspace, request.get("payload") or {})
    print(json.dumps({"success": True, **result}, ensure_ascii=False))


if __name__ == "__main__":
    try:
        main()
    except Exception as exc:
        print(json.dumps({"success": False, "error": str(exc)}, ensure_ascii=False))
        raise SystemExit(1)
